# -*- coding: utf-8 -*-
"""
Internal leadership demo (CTO): how the current delivery model works for the
client's 5 projects + SAFe ways-of-working + Project 5 deep dive.
Names are placeholders (editable). Output is a fully-editable .pptx.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ════════════════════════════ EDITABLE CONTENT ════════════════════════════════
CLIENT = "HDFC Bank"  # placeholder — replace later
PROJECTS = [
    dict(name="Project 1", onsite=0, offshore=2, buffer=0, roles=[("ATM Test Engineer", 2, 32)]),
    dict(name="Project 2", onsite=2, offshore=0, buffer=0, roles=[("ATM Solution Engg", 1, 90), ("ATM QA", 1, 89)]),
    dict(name="Project 3", onsite=0, offshore=1, buffer=0, roles=[("Tech Lead", 1, 32)]),
    dict(name="Project 4", onsite=0, offshore=6, buffer=1, roles=[("Lead Dev", 1, 32), ("Lead Dev", 1, 35), ("Lead QA", 1, 35), ("Sr. Dev", 2, 30), ("Sr. QA", 1, 32)]),
    dict(name="Project 5", onsite=0, offshore=10, buffer=1, roles=[("Sr Dev", 3, 28), ("Sr Dev", 2, 30), ("Tech Lead (Dev)", 1, 35), ("Sr QA", 1, 28), ("Lead QA", 2, 35), ("PO", 1, 37)]),
]
for p in PROJECTS:
    p["billable"] = sum(c for _, c, _ in p["roles"])
P5_TEAMS = [("Team 1", "HDFC FTE\n+ 2 Opus"), ("Team 2", "HDFC FTE\n+ 6 Opus + 1 buffer"), ("Team 3", "HDFC FTE\n+ 1 Opus")]
tot_billable = sum(p["billable"] for p in PROJECTS); tot_onsite = sum(p["onsite"] for p in PROJECTS)
tot_offshore = sum(p["offshore"] for p in PROJECTS); tot_buffer = sum(p["buffer"] for p in PROJECTS)
tot_opus = tot_billable + tot_buffer

# ── refined palette ───────────────────────────────────────────────────────────
NAVY=RGBColor(0x14,0x31,0x5A); BLUE=RGBColor(0x2F,0x6F,0xB0); LBLUE=RGBColor(0xEA,0xF2,0xFB)
SKY=RGBColor(0xBF,0xD6,0xEF); GREEN=RGBColor(0x2E,0x7D,0x32); GREENL=RGBColor(0xE7,0xF1,0xE8)
AMBER=RGBColor(0xC8,0x86,0x0D); GREY=RGBColor(0x6B,0x74,0x80); LINE=RGBColor(0xD6,0xDE,0xE8)
INK=RGBColor(0x22,0x2B,0x38); WHITE=RGBColor(0xFF,0xFF,0xFF); BG=RGBColor(0xF7,0xF9,0xFC)
FONT="Calibri"; FONT_H="Calibri"
PAGE=[0]

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
RR=MSO_SHAPE.ROUNDED_RECTANGLE

def slide(bg=None):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    if bg is not None:
        r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
        r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    return s
def rect(s,l,t,w,h,color,line=None,shape=MSO_SHAPE.RECTANGLE,lw=1.0,radius=None):
    sp=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h))
    if color is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False
    if radius is not None and shape==RR:
        try: sp.adjustments[0]=radius
        except Exception: pass
    return sp
def text(s,l,t,w,h,runs,*,size=14,color=INK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,italic=False,ls=1.0,font=FONT):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    items=runs if isinstance(runs,list) else [(runs,{})]
    for i,(txt,ov) in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=ov.get("align",align); p.line_spacing=ov.get("ls",ls)
        if "space_after" in ov: p.space_after=Pt(ov["space_after"])
        for j,seg in enumerate(txt if isinstance(txt,list) else [(txt,{})]):
            st,so=seg if isinstance(seg,tuple) else (seg,{})
            r=p.add_run(); r.text=st; f=r.font
            f.size=Pt(so.get("size",ov.get("size",size))); f.bold=so.get("bold",ov.get("bold",bold))
            f.italic=so.get("italic",ov.get("italic",italic)); f.name=ov.get("font",font)
            f.color.rgb=so.get("color",ov.get("color",color))
    return tb
def bullets(s,l,t,w,h,items,*,size=14.5,color=INK,gap=9):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=0; tf.margin_top=0
    for i,it in enumerate(items):
        head,body=(it if isinstance(it,tuple) else (it,None))
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap); p.line_spacing=1.05
        d=p.add_run(); d.text="—  "; d.font.size=Pt(size); d.font.name=FONT; d.font.bold=True; d.font.color.rgb=BLUE
        r=p.add_run(); r.text=head; r.font.size=Pt(size); r.font.name=FONT; r.font.bold=True; r.font.color.rgb=INK
        if body:
            r2=p.add_run(); r2.text="  "+body; r2.font.size=Pt(size); r2.font.name=FONT; r2.font.color.rgb=GREY
    return tb

def footer(s):
    text(s,0.6,7.12,9,0.28,f"{CLIENT}  ·  Delivery Model  ·  Internal leadership briefing (confidential)",size=8.5,color=GREY)
    text(s,12.0,7.12,0.85,0.28,str(PAGE[0]),size=8.5,color=GREY,align=PP_ALIGN.RIGHT)
def title_bar(s,title,message):
    PAGE[0]+=1
    rect(s,0,0,13.333,1.02,NAVY); rect(s,0,1.02,13.333,0.05,BLUE)
    text(s,0.62,0.15,11,0.24,"OPUS  ·  DELIVERY MODEL",size=10,color=SKY,bold=True)
    text(s,0.62,0.42,12,0.55,title,size=23,color=WHITE,bold=True,font=FONT_H)
    rect(s,0,1.07,13.333,0.5,LBLUE)
    text(s,0.62,1.13,12.1,0.4,[([("Key message:  ",{"bold":True,"color":BLUE}),(message,{"color":NAVY})],{})],size=12.5,anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
def kpi(s,l,t,w,h,label,value,sub,fill=NAVY):
    rect(s,l,t,w,h,fill,shape=RR,radius=0.05)
    text(s,l+0.18,t+0.10,w-0.36,0.24,label.upper(),size=10,color=SKY,bold=True)
    text(s,l+0.18,t+0.36,w-0.36,max(0.3,h-0.62),value,size=(22 if h>=1.05 else 18),color=WHITE,bold=True)
    text(s,l+0.18,t+h-0.26,w-0.36,0.24,sub,size=9.5,color=SKY)
def card(s,l,t,w,h,title,body_runs,accent=BLUE,tsize=13):
    rect(s,l,t,w,h,WHITE,line=LINE,lw=1.0)
    rect(s,l,t,0.085,h,accent)
    text(s,l+0.28,t+0.14,w-0.42,0.34,title,size=tsize,color=NAVY,bold=True)
    if body_runs is not None:
        text(s,l+0.28,t+0.52,w-0.46,h-0.62,body_runs,size=11.5,color=GREY,ls=1.08)
def fbox(s,l,t,w,h,title,sub,fill,tcolor=WHITE):
    rect(s,l,t,w,h,fill,shape=RR,radius=0.08)
    if sub:
        text(s,l+0.1,t+0.12,w-0.2,0.4,title,size=12.5,color=tcolor,bold=True,align=PP_ALIGN.CENTER)
        text(s,l+0.1,t+0.5,w-0.2,h-0.55,[(ln,{"align":PP_ALIGN.CENTER}) for ln in sub.split("\n")],size=10.5,color=tcolor,align=PP_ALIGN.CENTER)
    else:
        text(s,l+0.1,t,w-0.2,h,title,size=12.5,color=tcolor,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
def vline(s,x,t,h,color=LINE): rect(s,x,t,0.022,h,color)
def hline(s,l,t,w,color=LINE): rect(s,l,t,w,0.022,color)
def chev(s,l,t,w,h,label,fill):
    sp=rect(s,l,t,w,h,fill,shape=MSO_SHAPE.CHEVRON)
    text(s,l+0.18,t,w-0.3,h,label,size=12.5,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    return sp
def table(s,l,t,w,headers,rows,*,col_w=None,fs=11.5,hfs=11.5,row_h=0.34,head_h=0.42,highlight_rows=None,header_fill=NAVY):
    nrows=len(rows)+1; ncols=len(headers)
    gt=s.shapes.add_table(nrows,ncols,Inches(l),Inches(t),Inches(w),Inches(head_h+row_h*len(rows))).table
    if col_w:
        for i,cw in enumerate(col_w): gt.columns[i].width=Inches(cw)
    gt.first_row=False; gt.horz_banding=False; gt.rows[0].height=Inches(head_h)
    for j,htx in enumerate(headers):
        c=gt.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb=header_fill; c.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=c.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
        r=p.add_run(); r.text=htx; r.font.size=Pt(hfs); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=FONT
        c.margin_left=Inches(0.1); c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
    highlight_rows=highlight_rows or set()
    for i,row in enumerate(rows):
        gt.rows[i+1].height=Inches(row_h); hl=i in highlight_rows
        for j,val in enumerate(row):
            c=gt.cell(i+1,j); c.fill.solid(); c.fill.fore_color.rgb=GREENL if hl else (BG if i%2 else WHITE)
            c.vertical_anchor=MSO_ANCHOR.MIDDLE
            p=c.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
            r=p.add_run(); r.text=str(val); r.font.size=Pt(fs); r.font.name=FONT; r.font.bold=hl or j==0
            r.font.color.rgb=NAVY if (hl or j==0) else INK
            c.margin_left=Inches(0.1); c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
    return gt
def stacked(s,l,t,w,h,cats,series,*,title=None,colors=None):
    cd=CategoryChartData(); cd.categories=cats
    for name,vals in series: cd.add_series(name,vals)
    gf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED,Inches(l),Inches(t),Inches(w),Inches(h),cd)
    ch=gf.chart; ch.has_title=bool(title)
    if title:
        ch.chart_title.text_frame.text=title
        rr=ch.chart_title.text_frame.paragraphs[0].runs[0]; rr.font.size=Pt(12.5); rr.font.bold=True; rr.font.color.rgb=NAVY
    ch.has_legend=True; ch.legend.position=XL_LEGEND_POSITION.BOTTOM; ch.legend.include_in_layout=False; ch.legend.font.size=Pt(11)
    plot=ch.plots[0]; plot.gap_width=70; plot.has_data_labels=True
    plot.data_labels.font.size=Pt(10); plot.data_labels.font.bold=True; plot.data_labels.font.color.rgb=WHITE
    cset=colors or [BLUE,NAVY,AMBER]
    for i,sr in enumerate(ch.series):
        sr.format.fill.solid(); sr.format.fill.fore_color.rgb=cset[i%len(cset)]
    try: ch.value_axis.tick_labels.font.size=Pt(10); ch.category_axis.tick_labels.font.size=Pt(11)
    except Exception: pass
    ch.value_axis.has_major_gridlines=False
    return ch
def divider(num,title,subtitle):
    PAGE[0]+=1
    s=slide(NAVY); rect(s,0,3.62,13.333,0.05,BLUE)
    text(s,0.9,2.35,2.2,1.3,num,size=80,color=BLUE,bold=True)
    text(s,3.0,2.55,9.4,0.9,title,size=34,color=WHITE,bold=True,font=FONT_H)
    text(s,3.0,3.85,9.4,0.6,subtitle,size=15,color=SKY)
    text(s,0.9,7.0,9,0.3,f"{CLIENT}  ·  Delivery Model",size=9,color=GREY)
    return s

# ═══════════════════════════════ COVER ════════════════════════════════════════
s=slide(NAVY)
rect(s,0,0,0.28,7.5,BLUE)
rect(s,0.9,4.28,7.4,0.05,BLUE)
text(s,0.95,0.85,9,0.4,"OPUS  ·  DELIVERY EXCELLENCE",size=14,color=SKY,bold=True)
text(s,0.95,2.25,11.4,0.4,"INTERNAL LEADERSHIP BRIEFING  ·  CTO",size=13,color=SKY,bold=True)
text(s,0.93,2.75,11.5,1.5,f"{CLIENT} — How Our\nDelivery Model Works",size=40,color=WHITE,bold=True,ls=1.04,font=FONT_H)
text(s,0.95,4.5,11.4,0.6,"Portfolio  ·  Resourcing  ·  SAFe ways-of-working  ·  Project 5 deep dive",size=17,color=SKY)
text(s,0.95,6.75,11.4,0.4,"Confidential — Opus internal. Client and project names are placeholders.",size=10.5,color=GREY,italic=True)

# ═══════════════════════ ENGAGEMENT AT A GLANCE ══════════════════════════════
s=slide(BG); title_bar(s,"Engagement at a glance",
    f"Opus runs 5 projects for {CLIENT} — offshore-led, and now driving Project 5 toward production.")
kpi(s,0.62,1.8,2.95,1.35,"Projects",f"{len(PROJECTS)}","under one client",NAVY)
kpi(s,3.72,1.8,2.95,1.35,"Opus resources",f"{tot_opus}",f"{tot_billable} billable + {tot_buffer} buffer",BLUE)
kpi(s,6.82,1.8,2.95,1.35,"Delivery mix",f"{tot_offshore} : {tot_onsite}","offshore : onsite",NAVY)
kpi(s,9.92,1.8,2.8,1.35,"Framework","SAFe","4 PIs / year",GREEN)
card(s,0.62,3.45,5.95,1.55,"What we deliver",
     f"Five projects spanning ATM engineering & QA, development leads, and a full agile product team on Project 5 — staffed predominantly offshore ({tot_offshore} of {tot_billable} billable).",accent=BLUE)
card(s,6.78,3.45,5.95,1.55,"How the client operates",
     f"{CLIENT} runs the Scaled Agile Framework: every team plugs into a shared rhythm of 4 Program Increments a year and commits scope at PI planning.",accent=NAVY)
card(s,0.62,5.2,12.11,1.55,"Why this briefing",
     "Project 5 is our flagship — a three-team agile train inside a ~25-person ecosystem, with three customers heading to production (two in Q3 2026, a very large one in Q1 2027). This deck walks the model end-to-end.",accent=GREEN)

# ═══════════════════════ DIVIDER 1 ═══════════════════════════════════════════
divider("01","The Portfolio","Five projects, the people on them, and how they are deployed.")

# ═══════════════════════ PORTFOLIO SNAPSHOT ══════════════════════════════════
s=slide(BG); title_bar(s,"The portfolio",
    f"Across 5 projects we field {tot_opus} resources — {tot_offshore} offshore, {tot_onsite} onsite, {tot_buffer} on buffer.")
rows=[[p["name"],p["billable"],p["onsite"],p["offshore"],p["buffer"],p["billable"]+p["buffer"],"; ".join(f"{r} x{c} (${rt})" for r,c,rt in p["roles"])] for p in PROJECTS]
rows.append(["Total",tot_billable,tot_onsite,tot_offshore,tot_buffer,tot_opus,"—"])
table(s,0.62,1.85,12.1,["Project","Billable","Onsite","Offshore","Buffer","Total","Roles  (rate/hr)"],
      rows,col_w=[1.3,1.0,0.92,1.02,0.92,0.86,6.08],fs=11,hfs=11,row_h=0.66,highlight_rows={len(rows)-1})
text(s,0.62,6.4,11.0,0.4,"Billable = onsite + offshore.  Buffer resources cover ramp-up, leave and surge, and are not assigned to a billed feature.",size=10.5,color=GREY,italic=True)

# ═══════════════════════ RESOURCE DISTRIBUTION ═══════════════════════════════
s=slide(BG); title_bar(s,"Resource distribution",
    "An offshore-led model keeps cost efficient; a small onsite + buffer footprint protects coverage.")
rect(s,0.62,1.85,7.7,4.95,WHITE,line=LINE)
stacked(s,0.72,1.95,7.5,4.75,[p["name"] for p in PROJECTS],
        [("Offshore",[p["offshore"] for p in PROJECTS]),("Onsite",[p["onsite"] for p in PROJECTS]),("Buffer",[p["buffer"] for p in PROJECTS])],
        title="Resources by project",colors=[BLUE,NAVY,AMBER])
kpi(s,8.55,1.85,4.18,1.18,"Offshore share",f"{round(tot_offshore/tot_billable*100)}%","of billable resources",BLUE)
kpi(s,8.55,3.15,4.18,1.18,"Onsite presence",f"{tot_onsite}","client-site (Project 2)",NAVY)
kpi(s,8.55,4.45,4.18,1.18,"Buffer pool",f"{tot_buffer}","resilience on P4 & P5",AMBER)
card(s,8.55,5.75,4.18,1.05,"Read",
     "Projects scale from lean 1–2 person pods to a 10-person agile team (Project 5).",accent=GREEN)

# ═══════════════════════ RESOURCING DETAIL ═══════════════════════════════════
s=slide(BG); title_bar(s,"Resourcing detail",
    "The two largest projects carry the development leads (P4) and the full agile team (P5).")
def role_rows(p):
    rr=[[r,c,f"${rt}"] for r,c,rt in p["roles"]]
    rr.append(["Billable total",p["billable"],"—"]); rr.append(["Buffer",p["buffer"],"—"]); return rr
text(s,0.62,1.8,5.8,0.32,"Project 4",size=15,bold=True,color=NAVY)
table(s,0.62,2.18,5.7,["Role","Count","Rate/hr"],role_rows(PROJECTS[3]),col_w=[3.3,1.2,1.2],fs=12,row_h=0.42,highlight_rows={5,6})
text(s,6.95,1.8,5.8,0.32,[([("Project 5",{"bold":True,"size":15,"color":NAVY}),("   — flagship",{"size":12,"color":GREY,"italic":True})],{})])
table(s,6.95,2.18,5.7,["Role","Count","Rate/hr"],role_rows(PROJECTS[4]),col_w=[3.3,1.2,1.2],fs=12,row_h=0.42,highlight_rows={6,7})
card(s,0.62,6.0,12.11,0.85,"Projects 1–3 are lean pods",
     "ATM Test Engineers (×2, offshore) · an onsite ATM Solution Engg + ATM QA · a single Tech Lead — small, focused, and cost-light.",accent=BLUE)

# ═══════════════════════ DIVIDER 2 ═══════════════════════════════════════════
divider("02","How the Client Works","The Scaled Agile Framework — the rhythm every team delivers to.")

# ═══════════════════════ SAFe + PI CADENCE ═══════════════════════════════════
s=slide(BG); title_bar(s,"Ways of working — SAFe",
    "The year runs as 4 Program Increments; every team plans together and commits scope per PI.")
text(s,0.62,1.75,12,0.3,"ONE YEAR  =  4 PROGRAM INCREMENTS (PI)",size=12,bold=True,color=BLUE)
hline(s,0.9,3.35,11.5,LINE)
for i in range(4):
    x=0.9+i*2.92
    fbox(s,x,1.95,2.55,1.15,f"PI {i+1}","5–7 sprints",NAVY if i%2==0 else BLUE)
    # node on timeline
    rect(s,x+1.2,3.24,0.22,0.22,BLUE,shape=MSO_SHAPE.OVAL)
    text(s,x,3.55,2.55,0.4,"~10–14 weeks",size=10,color=GREY,align=PP_ALIGN.CENTER)
    # sprint ticks
    for k in range(6):
        rect(s,x+0.18+k*0.38,2.62,0.26,0.16,WHITE if i%2==0 else LBLUE,shape=RR,radius=0.3)
card(s,0.62,4.35,5.95,2.05,"The cadence",
     [("Each PI is 5–7 two-week sprints — roughly 10–14 weeks.",{"space_after":5}),
      ("All teams plan the PI together (PI planning), surface dependencies, and commit a set of features.",{"space_after":5}),
      ("Scope is sized in story points (SP); past velocity informs each new commitment.",{})],accent=BLUE)
card(s,6.78,4.35,5.95,2.05,"Worked example (illustrative)",
     [("A 6-sprint PI = 12 weeks.",{"bold":True,"color":INK,"space_after":5}),
      ("A team of 8 at ~8 SP/person/sprint ≈ 64 SP per sprint → ~380 SP across the PI.",{"color":GREY,"space_after":5}),
      ("That capacity is split between committed features and production-support work.",{"color":GREY})],accent=GREEN)

# ═══════════════════════ HOW A PI RUNS ═══════════════════════════════════════
s=slide(BG); title_bar(s,"How a PI runs",
    "Plan together, commit in story points, execute in sprints, and run support alongside.")
steps=[("PI Planning","All teams plan the PI together; dependencies and risks are surfaced up front.",NAVY),
       ("Commit features","Each team commits a set of features, sized in story points, for the PI.",BLUE),
       ("Execute sprints","5–7 two-week sprints build the committed scope, tracked sprint by sprint.",NAVY),
       ("Run + support","Teams balance committed delivery with incoming production-support issues.",BLUE)]
cw=3.0
for i,(t,d,col) in enumerate(steps):
    x=0.7+i*3.05
    chev(s,x,1.95,cw,0.95,f"{i+1}.  {t}",col)
    text(s,x+0.2,3.1,cw-0.35,1.5,d,size=11.5,color=GREY,ls=1.1)
card(s,0.62,4.85,12.11,1.7,"What leadership should take away",
     [([("Capacity is predictable.",{"bold":True,"color":INK}),("  Each team commits a story-point target at PI planning, so a PI's output is known up front and tracked every sprint.",{"color":GREY})],{"space_after":7}),
      ([("Shared services support the train.",{"bold":True,"color":INK}),("  Architecture, DB, DevOps and documentation back every team across all projects.",{"color":GREY})],{})],accent=GREEN)

# ═══════════════════════ DIVIDER 3 ═══════════════════════════════════════════
divider("03","Project 5 — Deep Dive","Our flagship: a three-team agile train heading to production.")

# ═══════════════════════ P5 TEAM STRUCTURE ═══════════════════════════════════
s=slide(BG); title_bar(s,"Project 5 — team structure",
    "Opus' 10 resources work across 3 teams inside a ~25-person ecosystem, backed by shared services.")
# top card
fbox(s,3.97,1.8,5.4,0.95,"Project 5",f"~25 people total · HDFC FTE + Opus (10 +1 buffer) + other vendors",NAVY)
# bus
vline(s,6.62,2.75,0.45,LINE); hline(s,1.9,3.2,9.5,LINE)
teamcols=[("Team 1","HDFC FTE\n+ 2 Opus",BLUE),("Team 2","HDFC FTE\n+ 6 Opus + 1 buffer",BLUE),("Team 3","HDFC FTE\n+ 1 Opus",BLUE),("Shared Services","Architect · DB Admin\nDevOps · Tech Writer",GREEN)]
xs=[0.62,3.66,6.7,9.74]
for (t,sub,col),x in zip(teamcols,xs):
    vline(s,x+1.45,3.2,0.42,LINE)
    fbox(s,x,3.62,2.95,1.5,t,sub,col)
card(s,0.62,5.45,12.11,1.35,"How to read this",
     [([("Opus is one part of a larger train.",{"bold":True,"color":INK}),("  Of ~25 people, 10 Opus (+1 buffer) sit alongside client full-time employees and other vendors, split 2 / 6 / 1 across three teams.",{"color":GREY})],{"space_after":7}),
      ([("Shared services",{"bold":True,"color":INK}),(" — architecture, DB, DevOps and technical writing — support all three teams.",{"color":GREY})],{})],accent=BLUE)

# ═══════════════════════ P5 COMPOSITION ══════════════════════════════════════
s=slide(BG); title_bar(s,"Project 5 — Opus composition",
    "The billable role mix, and how those resources are deployed across the three teams.")
text(s,0.62,1.8,5.8,0.32,"Billable role mix",size=14,bold=True,color=NAVY)
table(s,0.62,2.18,5.8,["Role","Count","Rate/hr"],role_rows(PROJECTS[4]),col_w=[3.4,1.2,1.2],fs=12,row_h=0.42,highlight_rows={6,7})
text(s,6.95,1.8,5.8,0.32,"Deployment across teams",size=14,bold=True,color=NAVY)
table(s,6.95,2.18,5.8,["Team","Opus resources"],[["Team 1","2"],["Team 2","6  (+1 buffer)"],["Team 3","1"],["Across 3 teams","10  (+1 buffer)"]],
      col_w=[3.0,2.8],fs=12.5,row_h=0.52,highlight_rows={3})
card(s,6.95,4.85,5.8,1.95,"Note",
     "The role table is the billable composition. The buffer resource absorbs leave, ramp-up and surge so committed velocity is protected when someone is unavailable.",accent=AMBER)

# ═══════════════════════ P5 GO-LIVE ROADMAP ══════════════════════════════════
s=slide(BG); title_bar(s,"Project 5 — go-live roadmap",
    "Three customers head to production: two in Q3 2026 and a very large one in Q1 2027.")
hline(s,1.0,4.0,11.3,LINE)
miles=[(2.4,"Q3 2026","Customer A","10-bank integration",GREEN),
       (6.0,"Q3 2026","Customer B","Production go-live",BLUE),
       (10.2,"Q1 2027","Customer C","Very large customer",NAVY)]
for x,q,name,note,col in miles:
    text(s,x-1.3,2.05,2.8,0.34,q,size=13,bold=True,color=col,align=PP_ALIGN.CENTER)
    rect(s,x-1.25,2.42,2.7,1.2,WHITE,line=col,lw=1.25)
    rect(s,x-1.25,2.42,2.7,0.1,col)
    text(s,x-1.15,2.62,2.5,0.4,name,size=13,bold=True,color=NAVY,align=PP_ALIGN.CENTER)
    text(s,x-1.15,3.05,2.5,0.5,note,size=11,color=GREY,align=PP_ALIGN.CENTER)
    rect(s,x-0.11,3.88,0.24,0.24,col,shape=MSO_SHAPE.OVAL)
    vline(s,x,3.62,0.3,col)
card(s,0.62,4.7,12.11,1.85,"What this means",
     [("Two go-lives land in Q3 2026; a very large customer follows in Q1 2027.",{"bold":True,"color":INK,"space_after":5}),
      ("Customer A integrates with 10 banks — we deploy the first bank in Q3 2026, then roll out the remaining nine once Bank 1 is proven (next slide).",{"color":GREY,"space_after":5}),
      ("The PI cadence absorbs go-live hardening and hypercare alongside committed features.",{"color":GREY})],accent=GREEN)

# ═══════════════════════ 10-BANK ROLLOUT ═════════════════════════════════════
s=slide(BG); title_bar(s,"Customer A — phased 10-bank rollout",
    "We de-risk by proving the integration on one bank, then scaling the proven pattern to the rest.")
chev(s,0.7,1.9,3.7,0.95,"1.  Bank 1 — deploy Q3 2026",GREEN)
chev(s,4.75,1.9,3.7,0.95,"2.  Validate & stabilize",AMBER)
chev(s,8.8,1.9,3.7,0.95,"3.  Banks 2–10 — phased",BLUE)
text(s,0.9,3.05,3.5,0.7,"First production go-live; prove the end-to-end integration.",size=11,color=GREY,align=PP_ALIGN.CENTER)
text(s,4.95,3.05,3.5,0.7,"Hypercare, confirm the integration, capture learnings.",size=11,color=GREY,align=PP_ALIGN.CENTER)
text(s,9.0,3.05,3.5,0.7,"Roll out the remaining nine banks, reusing the proven pattern.",size=11,color=GREY,align=PP_ALIGN.CENTER)
text(s,0.7,3.95,12,0.3,"10-BANK INTEGRATION",size=11,bold=True,color=NAVY)
for i in range(10):
    x=0.7+i*1.205
    on=(i==0)
    rect(s,x,4.32,1.02,0.72,GREEN if on else WHITE,line=None if on else LINE,shape=RR,radius=0.15)
    text(s,x,4.45,1.02,0.45,f"Bank {i+1}",size=10,color=WHITE if on else GREY,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
text(s,0.7,5.1,12,0.3,[([("■  ",{"color":GREEN}),("Live in Q3 2026          ",{"color":GREY}),("□  ",{"color":GREY}),("Phased after Bank 1 is proven",{"color":GREY})],{})],size=10.5)
card(s,0.62,5.55,12.11,1.2,"Why phased",
     "Proving the integration end-to-end on one bank — then hardening it — lets us scale a known-good pattern to the remaining nine, protecting both the customer relationship and our delivery commitments.",accent=GREEN)

# ═══════════════════════ KEY TAKEAWAYS ═══════════════════════════════════════
s=slide(BG); title_bar(s,"Summary","Five projects, one SAFe rhythm, and a flagship heading to production.")
bullets(s,0.7,1.9,12.0,4.6,[
    (f"Footprint:", f"5 projects, {tot_opus} Opus resources ({tot_billable} billable + {tot_buffer} buffer), delivered offshore-led with a small onsite + buffer presence."),
    ("Ways of working:", "the client runs SAFe — 4 PIs/year, 5–7 two-week sprints per PI — and every team commits a story-point target at PI planning."),
    ("Flagship:", "Project 5 is a three-team agile train (2 / 6 / 1 Opus + buffer) inside a ~25-person ecosystem, backed by shared services."),
    ("Production:", "three customers go live — two in Q3 2026 and a very large one in Q1 2027."),
    ("De-risking:", "Customer A's 10-bank integration rolls out phased — Bank 1 first, then the remaining nine once proven."),
],size=15,gap=14)
rect(s,0.62,6.35,12.11,0.55,LBLUE)
text(s,0.8,6.43,12,0.4,"All names are placeholders and every slide is fully editable in PowerPoint.",size=11.5,color=NAVY,italic=True,anchor=MSO_ANCHOR.MIDDLE)

def save():
    name="Opus_HDFC_Delivery_Model_Demo.pptx"
    try: prs.save(name); print("Saved",name)
    except PermissionError:
        alt=name.replace(".pptx","_UPDATED.pptx"); prs.save(alt); print(name,"open - saved",alt)
save()
print("Slides:",len(prs.slides._sldIdLst))
