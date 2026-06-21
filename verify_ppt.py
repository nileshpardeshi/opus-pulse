# -*- coding: utf-8 -*-
from pptx import Presentation
import os

def collect_text(prs):
    out = []
    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.has_text_frame:
                out.append(sh.text_frame.text)
            if sh.has_table:
                for row in sh.table.rows:
                    for c in row.cells:
                        out.append(c.text)
            if sh.has_chart:
                ch = sh.chart
                try:
                    if ch.has_title: out.append(ch.chart_title.text_frame.text)
                except Exception: pass
                for s in ch.series: out.append(str(s.name))
                try:
                    out += [str(c) for c in ch.plots[0].categories]
                except Exception: pass
    return "\n".join(out)

def stats(path):
    prs = Presentation(path)
    charts = tables = 0
    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.has_chart: charts += 1
            if sh.has_table: tables += 1
    return len(prs.slides.__iter__.__self__._sldIdLst), prs, charts, tables

for f in ["Opus_Pulse_Leadership_Internal.pptx", "Opus_Pulse_Client_Proposal.pptx"]:
    prs = Presentation(f)
    n = len(prs.slides._sldIdLst)
    charts = sum(1 for sl in prs.slides for sh in sl.shapes if sh.has_chart)
    tables = sum(1 for sl in prs.slides for sh in sl.shapes if sh.has_table)
    print(f"{f}: slides={n} charts={charts} tables={tables} size={round(os.path.getsize(f)/1024,1)}KB")

# RBAC leak check on the CLIENT deck
client = Presentation("Opus_Pulse_Client_Proposal.pptx")
txt = collect_text(client).lower()
forbidden = ["margin", "profit", "loaded cost", "opus cost", "28,900", "24,492",
             "45.3", "efficiency gain", "win-win", "230,9", "12,458", "12,444"]
hits = [w for w in forbidden if w in txt]
print("CLIENT deck forbidden-term hits (must be empty):", hits)
print("CLIENT deck OK" if not hits else "LEAK DETECTED -> fix")
