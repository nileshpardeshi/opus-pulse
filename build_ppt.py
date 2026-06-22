# -*- coding: utf-8 -*-
"""
Opus Pulse - TNM -> Outcome conversion: two professional decks.
  1) Opus_Pulse_Leadership_Internal.pptx  (full economics incl. cost & margin)
  2) Opus_Pulse_Client_Proposal.pptx       (price comparison + client benefits;
                                             NO internal cost/margin)
All numbers compute from the MODEL INPUTS block (single source of truth).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# ════════════════════════════════ MODEL INPUTS ════════════════════════════════
CLIENT = "Apex Bank"
ENGAGEMENT = "Payments Gateway"
R = 48147.0; RY = 577760.0; C = 28900.0
HEAD = 10; RESERVE = 0.15; EFF = 0.10; DISC = 0.05; SPRINTS = 2; MIN_SP = 120
CASE_A, CASE_B = 7, 8

tnm_profit=R-C; tnm_margin=tnm_profit/R
out_cost=C/(1+EFF); out_bill=R*(1-DISC); out_bill_y=out_bill*12
out_profit=out_bill-out_cost; out_margin=out_profit/out_bill
uplift=(out_margin-tnm_margin)*100
opus_extra_y=(out_profit-tnm_profit)*12
client_save_m=R-out_bill; client_save_y=client_save_m*12
eff_saving_m=C-out_cost; disc_given_m=R*DISC
def case(pp):
    gross=pp*HEAD; eff=gross*(1-RESERVE); msp=eff*SPRINTS; be=R/msp; rec=be*(1-DISC)
    return dict(pp=pp,gross=gross,eff=eff,msp=msp,be=be,rec=rec,meets=msp>=MIN_SP)
cA,cB=case(CASE_A),case(CASE_B)
def M(x): return f"${x:,.0f}"
def Mk(x): return f"${x/1000:,.1f}k"
def P(x): return f"{x*100:.1f}%"
def P0(x): return f"{x*100:.0f}%"

# ── refined palette ───────────────────────────────────────────────────────────
NAVY=RGBColor(0x14,0x31,0x5A); BLUE=RGBColor(0x2F,0x6F,0xB0); LBLUE=RGBColor(0xEA,0xF2,0xFB)
SKY=RGBColor(0xBF,0xD6,0xEF); GREEN=RGBColor(0x2E,0x7D,0x32); GREENL=RGBColor(0xE7,0xF1,0xE8)
AMBER=RGBColor(0xC8,0x86,0x0D); GREY=RGBColor(0x6B,0x74,0x80); LINE=RGBColor(0xD6,0xDE,0xE8)
INK=RGBColor(0x22,0x2B,0x38); WHITE=RGBColor(0xFF,0xFF,0xFF); BG=RGBColor(0xF7,0xF9,0xFC)
FONT="Calibri"; RR=MSO_SHAPE.ROUNDED_RECTANGLE
PAGE=[0]

def new_deck():
    p=Presentation(); p.slide_width=Inches(13.333); p.slide_height=Inches(7.5); return p
def slide(prs,bg=None):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    if bg is not None:
        r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
        r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    return s
def rect(s,l,t,w,h,color,line=None,shape=MSO_SHAPE.RECTANGLE,lw=1.0,radius=None):
    sp=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h))
    if color is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False
    if radius is not None and shape==RR:
        try: sp.adjustments[0]=radius
        except Exception: pass
    return sp
def text(s,l,t,w,h,runs,*,size=14,color=INK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,italic=False,ls=1.0,font=FONT):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    items=runs if isinstance(runs,list) else [(runs,{})]
    for i,(txt,ov) in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=ov.get("align",align); p.line_spacing=ov.get("ls",ls)
        if "space_after" in ov: p.space_after=Pt(ov["space_after"])
        for seg in (txt if isinstance(txt,list) else [(txt,{})]):
            st,so=seg if isinstance(seg,tuple) else (seg,{})
            r=p.add_run(); r.text=st; f=r.font
            f.size=Pt(so.get("size",ov.get("size",size))); f.bold=so.get("bold",ov.get("bold",bold))
            f.italic=so.get("italic",ov.get("italic",italic)); f.name=font; f.color.rgb=so.get("color",ov.get("color",color))
    return tb
def bullets(s,l,t,w,h,items,*,size=14.5,gap=9):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=0; tf.margin_top=0
    for i,it in enumerate(items):
        head,body=(it if isinstance(it,tuple) else (it,None))
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap); p.line_spacing=1.05
        d=p.add_run(); d.text="—  "; d.font.size=Pt(size); d.font.name=FONT; d.font.bold=True; d.font.color.rgb=BLUE
        r=p.add_run(); r.text=head; r.font.size=Pt(size); r.font.name=FONT; r.font.bold=True; r.font.color.rgb=INK
        if body:
            r2=p.add_run(); r2.text="  "+body; r2.font.size=Pt(size); r2.font.name=FONT; r2.font.color.rgb=GREY
    return tb
CONF=[""]
def footer(s):
    text(s,0.6,7.12,9,0.28,CONF[0],size=8.5,color=GREY)
    text(s,12.0,7.12,0.85,0.28,str(PAGE[0]),size=8.5,color=GREY,align=PP_ALIGN.RIGHT)
def title_bar(s,eyebrow,title,message):
    PAGE[0]+=1
    rect(s,0,0,13.333,1.02,NAVY); rect(s,0,1.02,13.333,0.05,BLUE)
    text(s,0.62,0.15,11,0.24,eyebrow.upper(),size=10,color=SKY,bold=True)
    text(s,0.62,0.42,12,0.55,title,size=23,color=WHITE,bold=True)
    rect(s,0,1.07,13.333,0.5,LBLUE)
    text(s,0.62,1.13,12.1,0.4,[([("Key message:  ",{"bold":True,"color":BLUE}),(message,{"color":NAVY})],{})],size=12.5,anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
def kpi(s,l,t,w,h,label,value,sub,fill=NAVY):
    rect(s,l,t,w,h,fill,shape=RR,radius=0.05)
    text(s,l+0.18,t+0.10,w-0.36,0.24,label.upper(),size=10,color=SKY,bold=True)
    text(s,l+0.18,t+0.36,w-0.36,max(0.3,h-0.62),value,size=(21 if h>=1.05 else 18),color=WHITE,bold=True)
    text(s,l+0.18,t+h-0.26,w-0.36,0.24,sub,size=9,color=SKY)
def card(s,l,t,w,h,title,body,accent=BLUE,tsize=13):
    rect(s,l,t,w,h,WHITE,line=LINE,lw=1.0); rect(s,l,t,0.085,h,accent)
    text(s,l+0.28,t+0.14,w-0.42,0.34,title,size=tsize,color=NAVY,bold=True)
    if body is not None: text(s,l+0.28,t+0.52,w-0.46,h-0.62,body,size=11.5,color=GREY,ls=1.08)
def fbox(s,l,t,w,h,title,sub,fill,tcolor=WHITE):
    rect(s,l,t,w,h,fill,shape=RR,radius=0.08)
    if sub:
        text(s,l+0.1,t+0.12,w-0.2,0.4,title,size=12.5,color=tcolor,bold=True,align=PP_ALIGN.CENTER)
        text(s,l+0.1,t+0.5,w-0.2,h-0.55,[(ln,{"align":PP_ALIGN.CENTER}) for ln in sub.split("\n")],size=10.5,color=tcolor,align=PP_ALIGN.CENTER)
    else:
        text(s,l+0.1,t,w-0.2,h,title,size=12.5,color=tcolor,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
def chev(s,l,t,w,h,label,fill):
    rect(s,l,t,w,h,fill,shape=MSO_SHAPE.CHEVRON)
    text(s,l+0.18,t,w-0.3,h,label,size=12,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
def chip(s,l,t,w,txt,color,fill):
    rect(s,l,t,w,0.46,fill,shape=RR,radius=0.5)
    text(s,l,t+0.05,w,0.34,txt,size=12,color=color,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
def table(s,l,t,w,headers,rows,*,col_w=None,fs=11.5,hfs=11.5,row_h=0.34,head_h=0.42,highlight_rows=None,header_fill=NAVY):
    nrows=len(rows)+1; ncols=len(headers)
    gt=s.shapes.add_table(nrows,ncols,Inches(l),Inches(t),Inches(w),Inches(head_h+row_h*len(rows))).table
    if col_w:
        for i,cw in enumerate(col_w): gt.columns[i].width=Inches(cw)
    gt.first_row=False; gt.horz_banding=False; gt.rows[0].height=Inches(head_h)
    for j,htx in enumerate(headers):
        c=gt.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb=header_fill; c.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=c.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
        r=p.add_run(); r.text=htx; r.font.size=Pt(hfs); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=FONT
        c.margin_left=Inches(0.1); c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
    highlight_rows=highlight_rows or set()
    for i,row in enumerate(rows):
        gt.rows[i+1].height=Inches(row_h); hl=i in highlight_rows
        for j,val in enumerate(row):
            c=gt.cell(i+1,j); c.fill.solid(); c.fill.fore_color.rgb=GREENL if hl else (BG if i%2 else WHITE)
            c.vertical_anchor=MSO_ANCHOR.MIDDLE
            p=c.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
            r=p.add_run(); r.text=str(val); r.font.size=Pt(fs); r.font.name=FONT; r.font.bold=hl or j==0
            r.font.color.rgb=NAVY if (hl or j==0) else INK
            c.margin_left=Inches(0.1); c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
    return gt
def clustered(s,l,t,w,h,cats,series,*,title=None,number_fmt='"$"#,##0',colors=None,pct=False,legend=True):
    cd=CategoryChartData(); cd.categories=cats
    for name,vals in series: cd.add_series(name,vals)
    gf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,Inches(l),Inches(t),Inches(w),Inches(h),cd)
    ch=gf.chart; ch.has_title=bool(title)
    if title:
        ch.chart_title.text_frame.text=title
        rr=ch.chart_title.text_frame.paragraphs[0].runs[0]; rr.font.size=Pt(12.5); rr.font.bold=True; rr.font.color.rgb=NAVY
    ch.has_legend=legend
    if legend:
        ch.legend.position=XL_LEGEND_POSITION.BOTTOM; ch.legend.include_in_layout=False; ch.legend.font.size=Pt(11)
    plot=ch.plots[0]; plot.gap_width=80; plot.has_data_labels=True
    plot.data_labels.number_format='0.0%' if pct else number_fmt; plot.data_labels.number_format_is_linked=False
    plot.data_labels.font.size=Pt(10); plot.data_labels.font.bold=True; plot.data_labels.position=XL_LABEL_POSITION.OUTSIDE_END
    cset=colors or [GREY,GREEN]
    for i,sr in enumerate(ch.series):
        sr.format.fill.solid(); sr.format.fill.fore_color.rgb=cset[i%len(cset)]
    try:
        ch.value_axis.has_major_gridlines=False; ch.value_axis.tick_labels.font.size=Pt(10); ch.category_axis.tick_labels.font.size=Pt(11)
    except Exception: pass
    return ch
def donut(s,l,t,w,h,title,data,colors):
    cd=CategoryChartData(); cd.categories=[d[0] for d in data]; cd.add_series("share",[d[1] for d in data])
    gf=s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT,Inches(l),Inches(t),Inches(w),Inches(h),cd)
    ch=gf.chart; ch.has_title=bool(title)
    if title:
        ch.chart_title.text_frame.text=title
        rr=ch.chart_title.text_frame.paragraphs[0].runs[0]; rr.font.size=Pt(12); rr.font.bold=True; rr.font.color.rgb=NAVY
    ch.has_legend=True; ch.legend.position=XL_LEGEND_POSITION.BOTTOM; ch.legend.include_in_layout=False; ch.legend.font.size=Pt(10.5)
    plot=ch.plots[0]; plot.has_data_labels=True
    dl=plot.data_labels; dl.number_format='0%'; dl.number_format_is_linked=False
    dl.show_percentage=True; dl.show_value=False; dl.font.size=Pt(11); dl.font.bold=True; dl.font.color.rgb=WHITE
    pts=ch.series[0].points
    for i,c in enumerate(colors):
        pts[i].format.fill.solid(); pts[i].format.fill.fore_color.rgb=c
    return ch
def divider(prs,num,title,subtitle):
    PAGE[0]+=1
    s=slide(prs,NAVY); rect(s,0,3.62,13.333,0.05,BLUE)
    text(s,0.9,2.3,2.4,1.3,num,size=80,color=BLUE,bold=True)
    text(s,3.05,2.5,9.4,0.9,title,size=33,color=WHITE,bold=True)
    text(s,3.05,3.8,9.4,0.6,subtitle,size=15,color=SKY)
    text(s,0.9,7.0,9,0.3,f"{CLIENT}  ·  {ENGAGEMENT}  ·  TNM -> Outcome",size=9,color=GREY)
    return s
def cover(prs,eyebrow,title,subtitle,tag):
    s=slide(prs,NAVY); rect(s,0,0,0.28,7.5,BLUE); rect(s,0.9,4.28,7.4,0.05,BLUE)
    text(s,0.95,0.85,9,0.4,"OPUS  ·  DELIVERY ECONOMICS",size=14,color=SKY,bold=True)
    text(s,0.95,2.2,11.4,0.4,eyebrow.upper(),size=13,color=SKY,bold=True)
    text(s,0.93,2.7,11.5,1.5,title,size=37,color=WHITE,bold=True,ls=1.04)
    text(s,0.95,4.5,11.4,0.7,subtitle,size=17,color=SKY)
    text(s,0.95,6.75,11.4,0.4,tag,size=10.5,color=GREY,italic=True)
    return s
def save_pptx(prs,name):
    try: prs.save(name); print("Saved",name)
    except PermissionError:
        alt=name.replace(".pptx","_UPDATED.pptx"); prs.save(alt); print(name,"open in PowerPoint - saved",alt)

# ════════════════════════════ DECK 1 — INTERNAL ══════════════════════════════
def build_internal():
    prs=new_deck(); CONF[0]="Confidential - Opus internal (contains cost & margin)"; global PAGE
    cover(prs,"Internal leadership briefing  ·  CTO",
          f"{CLIENT} — {ENGAGEMENT}\nTNM to Outcome Conversion",
          "The margin case for moving from Time-&-Material to Outcome-based delivery",
          "Confidential - Opus internal only. Contains cost & margin.")

    # 1 Executive summary
    s=slide(prs,BG); title_bar(s,"Executive summary","The opportunity",
        f"Convert {ENGAGEMENT} to Outcome: the client pays ~{P0(DISC)} less while our margin rises {P0(tnm_margin)} -> {P(out_margin)} and the efficiency we unlock becomes ours to keep.")
    kpi(s,0.62,1.8,2.95,1.3,"Margin today",P0(tnm_margin),f"{Mk(RY)}/yr billing",NAVY)
    kpi(s,3.72,1.8,2.95,1.3,"Outcome margin",P(out_margin),f"+{uplift:.1f} points",GREEN)
    kpi(s,6.82,1.8,2.95,1.3,"Opus extra profit",f"+{Mk(opus_extra_y)}/yr","plus the mix upside",GREEN)
    kpi(s,9.92,1.8,2.8,1.3,"Client saving",f"{Mk(client_save_y)}/yr",f"~{P0(DISC)} less",BLUE)
    card(s,0.62,3.4,3.85,3.25,"The ask",
         f"Move {ENGAGEMENT} from billing time at a frozen rate to a fixed price per accepted story point, anchored to today's run-rate.",accent=NAVY)
    card(s,4.64,3.4,3.85,3.25,"Why it works",
         "Under TNM, any efficiency we create just reduces billed hours and is handed to the client. Under Outcome the price is fixed - a leaner resource mix and trimmed overhead become OUR margin, and we share a small slice back.",accent=BLUE)
    card(s,8.66,3.4,4.07,3.25,"The recommendation",
         "Pilot the conversion on this engagement, commit ~120 SP/month, lock the price-per-point, and stand up the resource pyramid that funds the efficiency. Expand on success.",accent=GREEN)

    # ── DIVIDER 1 ──
    divider(prs,"01","The Problem","Why Time-&-Material quietly caps our margin.")

    # 2 Current engagement snapshot
    s=slide(prs,BG); title_bar(s,"Current engagement","Where we are today",
        f"A 10-person team bills {M(R)}/month at ~{P0(tnm_margin)} margin - healthy, but the rate is frozen.")
    table(s,0.62,1.8,7.3,["Role","Exp (yrs)","Count","Rate/hr"],
          [["Sr Dev","5 to 8","3","$28"],["Sr Dev","8 to 10","2","$30"],["Tech Lead (Dev)","9 to 12","1","$35"],
           ["Sr QA","5 to 8","1","$28"],["Lead QA","9 to 12","2","$35"],["PO","12 to 14","1","$37"],
           ["Total team","","10","$314/hr blended"]],
          col_w=[3.0,1.55,1.2,1.55],highlight_rows={6},fs=11.5,row_h=0.42)
    kpi(s,8.2,1.8,4.53,0.92,"Billable days / year","230","365 - 104 wknd - 10 hol - 21 leave",NAVY)
    kpi(s,8.2,2.85,4.53,0.92,"Monthly billing",M(R),"153.3 hrs/mo x $314 blended",BLUE)
    kpi(s,8.2,3.9,4.53,0.92,"Yearly billing",M(RY),"frozen at the SOW rate",BLUE)
    card(s,0.62,5.1,12.11,1.5,"How the numbers are built",
         [("Billable hours/yr = 230 days x 8 hrs = 1,840 hrs -> 153.3 hrs/month.",{"space_after":4}),
          (f"Monthly billing = (count x rate) x billable hours = $314/hr x 153.3 = {M(R)}.",{"space_after":4}),
          (f"Loaded cost (salaries, infra, licenses, travel, training) ~ {M(C)}/month -> {P0(tnm_margin)} margin.",{})],accent=BLUE)

    # 3 TNM margin math
    s=slide(prs,BG); title_bar(s,"Current economics","The TNM margin math",
        f"We bill {M(R)}/mo against ~{M(C)} cost - a healthy {P0(tnm_margin)} margin that erodes as cost rises and the rate stays flat.")
    table(s,0.62,1.85,5.9,["Line","Monthly","Yearly"],
          [["Billing (revenue)",M(R),M(RY)],["Loaded cost",M(C),M(C*12)],["Profit",M(tnm_profit),M(tnm_profit*12)],["Margin %",P(tnm_margin),P(tnm_margin)]],
          col_w=[2.6,1.65,1.65],highlight_rows={3},row_h=0.5,fs=13)
    rect(s,6.75,1.85,6.0,4.0,WHITE,line=LINE)
    clustered(s,6.9,1.95,5.7,3.55,["Revenue","Cost","Profit"],[("Monthly $",[R,C,tnm_profit])],title="TNM monthly economics",colors=[BLUE],legend=False)
    card(s,0.62,4.2,5.9,2.4,"The catch",
         [("Margin is healthy today - but frozen.",{"bold":True,"color":INK,"space_after":5}),
          ("Appraisals raise cost each year while the bill rate stays flat, so margin compresses over the SOW.",{"space_after":5}),
          ("And any productivity we gain just reduces billed hours - the saving goes to the client, not us.",{})],accent=AMBER)

    # 4 Why TNM caps us
    s=slide(prs,BG); title_bar(s,"The problem","Why TNM caps our upside",
        "Three structural limits: frozen revenue, punished efficiency, and unpriced delivery risk.")
    probs=[("1.  Revenue is frozen","Locked at the SOW rate. Cost rises every year (appraisals, attrition backfills), so margin silently compresses.",AMBER),
           ("2.  Efficiency is punished","Deliver the same scope with fewer / cheaper people and we simply bill fewer hours - the client keeps 100% of the saving.",AMBER),
           ("3.  Risk is unpriced","We carry attrition, ramp-up and rework, but TNM gives us no way to price for it.",AMBER)]
    for i,(t,d,col) in enumerate(probs):
        x=0.62+i*4.06
        card(s,x,1.9,3.85,3.0,t,d,accent=col,tsize=14)
    rect(s,0.62,5.2,12.11,1.35,NAVY,shape=RR,radius=0.04)
    text(s,0.9,5.32,11.6,1.1,[([("Outcome-based delivery flips all three:  ",{"bold":True,"color":WHITE}),("the price is fixed to scope, a leaner mix becomes our margin, and delivery risk is priced into the bid.",{"color":SKY})],{})],size=14,anchor=MSO_ANCHOR.MIDDLE)

    # ── DIVIDER 2 ──
    divider(prs,"02","The Outcome Model","Price the outcome, keep the efficiency.")

    # 5 Outcome mechanics + cases
    s=slide(prs,BG); title_bar(s,"Proposed model","Mechanics & capacity cases",
        f"Capacity -> velocity -> price per point; we test two cases ({CASE_A} vs {CASE_B} SP/person) with a {P0(RESERVE)} reserve.")
    flow=[("Per-person capacity",f"{CASE_A}-{CASE_B} SP / sprint x 10 people",NAVY),
          (f"Less {P0(RESERVE)} reserve","meetings, grooming, issues",BLUE),
          (f"x {SPRINTS} sprints","= deliverable SP / month",NAVY),
          ("Price per point","run-rate / SP, then share-back",GREEN)]
    for i,(t,d,col) in enumerate(flow):
        x=0.62+i*3.06
        chev(s,x,1.85,2.95,0.85,t,col)
        text(s,x+0.15,2.78,2.7,0.6,d,size=10.5,color=GREY,align=PP_ALIGN.CENTER)
    table(s,0.62,3.7,12.11,["Metric",f"Case 1  ({CASE_A} SP/person)",f"Case 2  ({CASE_B} SP/person)"],
          [["Gross team velocity / sprint",f"{cA['gross']:.0f} SP",f"{cB['gross']:.0f} SP"],
           [f"Less {P0(RESERVE)} reserve -> effective",f"{cA['eff']:.1f} SP",f"{cB['eff']:.0f} SP"],
           ["Deliverable SP / month",f"{cA['msp']:.0f} SP",f"{cB['msp']:.0f} SP"],
           [f"Meets {MIN_SP}-SP min invoice?",("Yes" if cA['meets'] else f"No ({cA['msp']:.0f}) - lower min to ~110"),("Yes" if cB['meets'] else "No")]],
          col_w=[4.3,3.9,3.9],highlight_rows={2},row_h=0.5,fs=12)
    text(s,0.62,6.5,12.11,0.5,"Both cases bill the same monthly total (price anchored to the run-rate) - velocity changes the rate per point and whether the minimum invoice is met.",size=11,color=GREY,italic=True)

    # 6 Pricing calc
    s=slide(prs,BG); title_bar(s,"Pricing","Outcome price calculation",
        f"Anchor to the TNM run-rate, share {P0(DISC)} with the client - and the hard floor is always respected.")
    table(s,0.62,1.85,7.5,["Step",f"Case 1 ({CASE_A} SP)",f"Case 2 ({CASE_B} SP)"],
          [["Deliverable SP / month",f"{cA['msp']:.0f}",f"{cB['msp']:.0f}"],
           ["Break-even price/SP (run-rate / SP)",M(cA['be']),M(cB['be'])],
           ["Client share-back",f"-{P0(DISC)}",f"-{P0(DISC)}"],
           ["Recommended price/SP",M(cA['rec']),M(cB['rec'])],
           ["Monthly bill",M(out_bill),M(out_bill)],
           ["Yearly bill",M(out_bill_y),M(out_bill_y)]],
          col_w=[3.7,1.9,1.9],highlight_rows={3,4},row_h=0.44,fs=12)
    kpi(s,8.35,1.85,4.38,0.95,"Recommended",f"{M(cB['rec'])}/SP",f"Case 2, {P0(DISC)} share-back",GREEN)
    kpi(s,8.35,2.95,4.38,0.95,"Monthly bill",M(out_bill),f"vs {M(R)} TNM",BLUE)
    kpi(s,8.35,4.05,4.38,0.95,"Outcome cost",M(out_cost),f"{Mk(C)} / (1 + {P0(EFF)} eff.)",NAVY)
    card(s,0.62,5.45,12.11,1.15,"Guardrail",
         f"The price is always clamped to >= cost / (1 - min margin) and >= the TNM run-rate. The {P0(DISC)} we give the client is funded by a realistic {P0(EFF)} delivery-efficiency we keep under Outcome (next slide).",accent=BLUE)

    # 7 Resource mix (where efficiency comes from)
    s=slide(prs,BG); title_bar(s,"Funding the gain - no new tools",f"Where the {P0(EFF)} efficiency comes from",
        "We already use AI tooling under TNM - the gain comes from what Outcome unlocks: a leaner resource mix.")
    levers=[("Resource pyramid","Price is fixed per point, so the client no longer pays per named resource. Use freshers for routine work with senior oversight; blended cost drops while the price holds.",BLUE),
            ("Talent rotation","As juniors gain skill (and cost), rotate them onto other billable projects; bring in fresh low-cost joiners.",NAVY),
            ("Trim overhead 20% -> 15%","Less ceremony / meeting time; the same team delivers more points per month.",BLUE),
            ("Leaner delivery","No gold-plating; tighter Definition-of-Done -> less rework, which is OUR cost under Outcome.",NAVY)]
    for i,(t,d,col) in enumerate(levers):
        x=0.62+(i%2)*6.07; y=1.85+(i//2)*1.5
        card(s,x,y,5.9,1.38,t,d,accent=col,tsize=12.5)
    table(s,0.62,4.95,7.3,["Per story point","Senior pod","Junior-heavy"],
          [["Cost / sprint","$4,000","$2,600  (-35%)"],["Velocity","8 SP","6 SP  (-25%)"],
           ["Cost / SP (after ~10% rework)","$500","$476"],["Net efficiency / point","-","~5-10%"]],
          col_w=[3.3,2.0,2.0],highlight_rows={3},row_h=0.4,fs=11.5)
    card(s,8.1,4.95,4.63,1.95,"The honest number",
         "A 35% salary gap nets only ~5-10% per point once slower velocity and rework are subtracted. Use a pyramid (juniors + senior oversight), never all-juniors, or you miss the commitment.",accent=AMBER)

    # ── DIVIDER 3 ──
    divider(prs,"03","The Payoff","A genuine win-win - and how we de-risk it.")

    # 8 Win-win economics
    s=slide(prs,BG); title_bar(s,"The payoff","TNM vs Outcome - the win-win",
        "The client pays less AND we earn more - because Outcome lets us keep efficiency TNM would give away.")
    rect(s,0.62,1.85,7.0,4.3,WHITE,line=LINE)
    clustered(s,0.72,1.95,6.8,3.9,["Client pays","Opus cost","Opus profit"],
              [("TNM (monthly)",[R,C,tnm_profit]),("Outcome (monthly)",[out_bill,out_cost,out_profit])],
              title="Monthly economics: TNM vs Outcome",colors=[GREY,GREEN])
    kpi(s,7.85,1.85,4.88,0.92,"Opus margin",f"{P0(tnm_margin)} -> {P(out_margin)}",f"+{uplift:.1f} points",GREEN)
    kpi(s,7.85,2.9,4.88,0.92,"Opus extra profit",f"+{M(opus_extra_y)}/yr","plus mix upside we retain",GREEN)
    kpi(s,7.85,3.95,4.88,0.92,"Client saving",f"{M(client_save_y)}/yr",f"client pays ~{P0(DISC)} less",BLUE)
    chip(s,7.85,5.0,4.88,"RESULT:   WIN  -  WIN",WHITE,GREEN)
    card(s,0.62,6.25,12.11,0.85,"Why it holds",
         f"Efficiency saving {M(eff_saving_m)}/mo is greater than the discount given {M(disc_given_m)}/mo - so the client pays less and our profit still rises. Dial the discount down (e.g. 4%) to widen Opus' share.",accent=GREEN)

    # 9 Risks
    s=slide(prs,BG); title_bar(s,"De-risking","Risks & how we control them",
        "The model is protected by pricing discipline and a senior-anchored delivery pyramid.")
    table(s,0.62,1.9,12.11,["Risk","Mitigation"],
          [["Juniors slow velocity / miss commit","Pyramid with senior oversight; anchor price to a conservative velocity; minimum-invoice floor."],
           ["Quality / rework on junior work","Acceptance gate; strong Definition-of-Done; rework is pre-acceptance (our cost, buffered)."],
           ["Rotation loses knowledge","Stagger rotations; document; keep a stable senior core on the pod."],
           ["Scope creep on fixed price","All changes via priced change requests - nothing absorbed for free."],
           ["Client rejects fixed price","Sequence: managed-capacity -> gain-share -> outcome to build trust first."]],
          col_w=[4.3,7.81],row_h=0.74,fs=12)

    # 10 Recommendation
    s=slide(prs,BG); title_bar(s,"Decision","Recommendation & next steps",
        "Pilot the conversion now; the mix saving is retained as it grows, and the client pays less.")
    bullets(s,0.7,1.95,12.0,3.9,[
        ("Approve a pilot conversion",f"of {ENGAGEMENT} to Outcome-based delivery."),
        ("Commit & price",f"~{cB['msp']:.0f} SP/month at {M(cB['rec'])}/SP (Case 2); lower the invoice minimum to ~110 if we staff to {CASE_A} SP/person."),
        ("Stand up the pyramid","fresher intake + senior oversight + rotation - the engine that funds the efficiency."),
        ("Tighten delivery","trim ceremonies to a 15% reserve and harden the Definition-of-Done to cut rework."),
        ("Lock the SOW addendum","fix the price-per-point and acceptance definition; review after two invoices, then expand."),
    ],size=15,gap=13)
    chip(s,0.7,6.25,7.4,f"Target:  {P0(tnm_margin)} -> {P(out_margin)} margin, same team",WHITE,GREEN)

    save_pptx(prs,"Opus_Pulse_Leadership_Internal.pptx")
    return PAGE[0]

# ════════════════════════════ DECK 2 — CLIENT ════════════════════════════════
def build_client():
    prs=new_deck(); CONF[0]=f"Prepared by Opus for {CLIENT}";
    eff_per_sp=R/cB['msp']
    cover(prs,"Partnership proposal",
          "A Better Way to Deliver\nFrom Time-&-Material to Outcome",
          "Pay for delivered outcomes - with price certainty and lower cost",
          f"Prepared by Opus for {CLIENT}")

    # 1 Today
    s=slide(prs,BG); title_bar(s,"Today","Where we are today",
        f"Your {ENGAGEMENT} team of 10 is billed for time - {M(R)}/month - regardless of output.")
    bullets(s,0.62,1.95,7.0,3.6,[
        ("Time-&-Material model","you are invoiced for hours x rate."),
        ("Current investment",f"{M(R)} / month ({M(RY)} / year) in the team."),
        ("It works","but it ties your spend to effort, not to results."),
    ],size=15,gap=13)
    kpi(s,7.85,1.95,4.88,1.15,"Current monthly spend",M(R),"Time & Material",NAVY)
    kpi(s,7.85,3.3,4.88,1.15,"Current yearly spend",M(RY),"team of 10",NAVY)
    kpi(s,7.85,4.65,4.88,1.15,"Typical delivery",f"~{cB['msp']:.0f} SP/mo","story points accepted",BLUE)

    # 2 The hidden cost of TNM
    s=slide(prs,BG); title_bar(s,"The challenge","The hidden cost of Time-&-Material",
        "Under TNM you pay for every hour - including the meetings, rework and onboarding behind the work.")
    rect(s,0.62,1.85,5.75,4.78,WHITE,line=LINE)
    donut(s,0.74,1.98,5.5,4.5,"Where your billed hours go (illustrative)",
          [("Feature delivery",65),("Meetings & grooming",20),("Rework & ramp-up",15)],[GREEN,GREY,AMBER])
    text(s,6.6,1.95,6.13,0.34,"All of it is on your invoice today",size=14,bold=True,color=NAVY)
    hc=[("Meetings & grooming","Ceremony and planning time you fund every sprint.",GREY),
        ("Rework & defects","Fixing issues is billed to you - not absorbed by the vendor.",AMBER),
        ("Ramp-up & churn","New joiners bill full rate while they get up to speed.",BLUE)]
    for i,(t,d,col) in enumerate(hc):
        card(s,6.6,2.42+i*1.42,6.13,1.28,t,d,accent=col,tsize=13)

    # 3 Worked example - the numbers
    s=slide(prs,BG); title_bar(s,"Worked example","What you actually pay",
        "Same delivery, two ways to pay - a typical month, and a single 80-point feature.")
    text(s,0.62,1.82,5.95,0.32,"A typical month  (~136 story points)",size=14,bold=True,color=NAVY)
    table(s,0.62,2.2,5.95,["","Today (TNM)","Outcome"],
          [["You pay",M(R),M(out_bill)],["For","~136 SP delivered","136 accepted SP"],["Per point",f"~{M(eff_per_sp)}",M(cB['rec'])]],
          col_w=[1.95,2.0,2.0],row_h=0.46,fs=12,highlight_rows={0})
    chip(s,0.62,4.3,5.95,f"You save ~{M(client_save_m)} / month  (~{P0(DISC)})",WHITE,GREEN)
    text(s,6.78,1.82,5.95,0.32,"A single feature  (80 story points)",size=14,bold=True,color=NAVY)
    table(s,6.78,2.2,5.95,["","Today (TNM)","Outcome"],
          [["Feature price",f"~{M(80*eff_per_sp)}",M(80*cB['rec'])],["Basis",f"80 x ~{M(eff_per_sp)}",f"80 x {M(cB['rec'])}"],["You pay","regardless","only on acceptance"]],
          col_w=[1.95,2.0,2.0],row_h=0.46,fs=12,highlight_rows={0})
    chip(s,6.78,4.3,5.95,f"You save {M(80*eff_per_sp-80*cB['rec'])} + risk on us",WHITE,GREEN)
    card(s,0.62,5.05,12.11,1.5,"How to read it",
         f"Today you pay for the hours behind a feature - meetings, rework and all - even if it slips. Under Outcome you pay a fixed {M(cB['rec'])} per accepted story point, and nothing until it is delivered and signed off.",accent=BLUE)

    # 4 Proposal
    s=slide(prs,BG); title_bar(s,"Our proposal","Outcome-based delivery",
        f"Pay a fixed price per accepted story point - and it costs ~{P0(DISC)} less than today.")
    bullets(s,0.62,1.95,7.1,3.7,[
        ("Fixed price per accepted point","not per hour."),
        ("A committed monthly volume","with an agreed acceptance definition."),
        ("You pay only for accepted work","rework before acceptance is on us."),
        ("Delivery risk moves to Opus","attrition, ramp-up and estimation are ours."),
    ],size=15,gap=12)
    kpi(s,7.95,1.95,4.78,1.15,"Proposed price",f"{M(cB['rec'])}/SP","per accepted story point",GREEN)
    kpi(s,7.95,3.3,4.78,1.15,"Monthly commitment",f"~{cB['msp']:.0f} SP",f"= {M(out_bill)} / month",BLUE)
    kpi(s,7.95,4.65,4.78,1.15,"You pay only for","accepted work","quality-gated",NAVY)

    # 4 What you pay
    s=slide(prs,BG); title_bar(s,"The numbers","What you pay - today vs proposed",
        f"You save ~{Mk(client_save_y)} a year (~{P0(DISC)}) - and only pay for results.")
    table(s,0.62,1.9,6.6,["","Time & Material","Outcome-based"],
          [["Basis","Hours x rate","Price per accepted SP"],
           ["Effective $ / delivered SP",f"~{M(eff_per_sp)}",M(cB['rec'])],
           ["Monthly invoice",M(R),M(out_bill)],
           ["Yearly invoice",M(RY),M(out_bill_y)],
           ["You pay for","All time","Accepted output only"]],
          col_w=[2.5,2.05,2.05],highlight_rows={2,3},row_h=0.5,fs=12)
    rect(s,7.45,1.9,5.28,3.1,WHITE,line=LINE)
    clustered(s,7.55,2.0,5.1,2.9,["Per month","Per year"],[("Time & Material",[R,RY]),("Outcome-based",[out_bill,out_bill_y])],
              title="Your spend: today vs proposed",colors=[GREY,GREEN])
    chip(s,0.62,5.55,6.6,f"YOUR SAVING:  ~{Mk(client_save_y)}/yr  (~{P0(DISC)})",WHITE,GREEN)
    card(s,7.45,5.15,5.28,1.5,"Read",
         f"Today {M(R)} buys ~{cB['msp']:.0f} delivered points (~{M(eff_per_sp)} each), paid whether or not they land. At {M(cB['rec'])} per accepted point you pay less - and only for results.",accent=BLUE)

    # 5 Beyond price
    s=slide(prs,BG); title_bar(s,"Beyond price","What you gain - beyond the saving",
        "Certainty, risk transfer and a delivery commitment - not just a lower number.")
    cards=[("Price certainty","A fixed rate per point. No surprises from rate revisions or slippage.",BLUE),
           ("Pay for results","Rework, meetings and ramp-up are no longer on your invoice.",GREEN),
           ("Risk transferred","Attrition, ramp-up and estimation risk sit with Opus.",NAVY),
           ("Throughput SLA","A committed monthly delivery volume - not an open-ended timesheet.",BLUE),
           ("Less overhead","You govern outcomes, not headcount and utilisation.",GREEN),
           ("Built-in quality","Acceptance gate - you sign off before it is billable.",NAVY)]
    for i,(t,d,col) in enumerate(cards):
        x=0.62+(i%3)*4.06; y=1.95+(i//3)*2.05
        card(s,x,y,3.85,1.85,t,d,accent=col,tsize=13.5)

    # 6 At a glance
    s=slide(prs,BG); title_bar(s,"At a glance","Time-&-Material vs Outcome-based",
        "The same delivery, repriced around results - lower cost and more certainty for you.")
    table(s,0.7,1.95,11.9,["Dimension","Time & Material","Outcome-based"],
          [["You pay for","Hours worked","Accepted story points"],
           ["Price certainty","Variable","Fixed per point"],
           ["Delivery risk","Yours","Opus"],
           ["Getting faster over time","Fewer billed hours","A lower price per point"],
           ["Rework / ramp-up","On your invoice","Absorbed by Opus"],
           ["Your effort","Manage the team","Approve the outcomes"],
           ["Annual cost",M(RY),M(out_bill_y)]],
          col_w=[3.4,4.25,4.25],highlight_rows={6},row_h=0.54,fs=12.5)

    # 7 Transition + close
    s=slide(prs,BG); title_bar(s,"Getting there","A low-risk transition",
        "Validate in parallel first, then switch the SOW - no leap of faith required.")
    bullets(s,0.7,1.95,12.0,3.1,[
        ("Step 1 - Agree the basics","story-point definition, acceptance criteria and monthly volume."),
        ("Step 2 - Run in parallel","one to two monthly cycles to validate velocity and pricing."),
        ("Step 3 - Switch the SOW","to a fixed price per accepted point."),
        ("Step 4 - Review & extend","quarterly reviews; extend across other workstreams."),
    ],size=15,gap=13)
    chip(s,0.7,5.25,6.8,f"You save ~{Mk(client_save_y)}/yr + gain certainty",WHITE,GREEN)
    text(s,0.7,5.95,12.0,1.0,[("Let's deliver outcomes, not hours.",{"size":22,"bold":True,"color":NAVY}),
                              ("We're ready to start the parallel run next sprint.",{"size":14,"color":GREY})])

    save_pptx(prs,"Opus_Pulse_Client_Proposal.pptx")
    return PAGE[0]

PAGE[0]=0; ni=build_internal()
PAGE[0]=0; nc=build_client()
print("Internal slides:",ni,"| Client slides:",nc)
print("Model: reserve",P0(RESERVE),"eff",P0(EFF),"discount",P0(DISC),
      "| client saves",M(client_save_y),"| opus extra",M(opus_extra_y),"| margin",P0(tnm_margin),"->",P(out_margin))
