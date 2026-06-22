# -*- coding: utf-8 -*-
"""Flags overlapping TEXT boxes within each slide (catches text-on-text overlap)."""
import sys
from pptx import Presentation
from pptx.util import Emu

def rects(slide):
    out = []
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            try:
                l, t, w, h = sh.left, sh.top, sh.width, sh.height
                if None in (l, t, w, h):
                    continue
                out.append((Emu(l).inches, Emu(t).inches, Emu(w).inches, Emu(h).inches, sh.text_frame.text.strip().replace("\n", " ")[:32]))
            except Exception:
                pass
    return out

def overlap(a, b):
    ax, ay, aw, ah, _ = a; bx, by, bw, bh, _ = b
    ix = max(0, min(ax + aw, bx + bw) - max(ax, bx))
    iy = max(0, min(ay + ah, by + bh) - max(ay, by))
    return ix, iy

for f in sys.argv[1:]:
    prs = Presentation(f); flags = 0
    for i, sl in enumerate(prs.slides, 1):
        rs = rects(sl)
        for m in range(len(rs)):
            for n in range(m + 1, len(rs)):
                ix, iy = overlap(rs[m], rs[n])
                # meaningful text-on-text overlap
                if iy > 0.10 and ix > 0.6:
                    flags += 1
                    print(f"  [{f}] slide {i}: '{rs[m][4]}'  <>  '{rs[n][4]}'  (overlap {ix:.2f}\" x {iy:.2f}\")")
    print(f"{f}: {'OK - no text overlaps' if flags==0 else str(flags)+' overlap(s)'}")
